import os
import re

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lex_rank import LexRankSummarizer


class PPTXItemPipeline(object):
    def __init__(self, file_name):
        self.__file_name = file_name

    @classmethod
    def from_crawler(cls, crawler):
        return cls(
            file_name=crawler.settings.get('PPTX_FILE_NAME'),
        )

    def open_spider(self, spider):
        self.__presentation = Presentation()

    def close_spider(self, spider):
        tag_data = sorted(
            (
                (k[5:], v) for (k, v) in
                spider.crawler.stats.get_stats().iteritems()
                if k.startswith('tags/')
            ),
            key=lambda x: x[1],
            reverse=True,
        )

        # Take only the most frequent tags.
        tag_data = tag_data[:15]

        import matplotlib.pyplot as plt
        plt.pie(
            x=[v for (k, v) in tag_data],
            labels=[k for (k, v) in tag_data],
            startangle=140,
        )
        plt.axis('equal')
        plt.savefig('cake.png', bbox_inches='tight')

        # Add the cake diagram.
        blank_slide_layout = self.__presentation.slide_layouts[5]
        slide = self.__presentation.slides.add_slide(blank_slide_layout)
        slide.shapes.title.text = 'Most Frequent Tags'
        slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(20)

        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
        chart_data = ChartData()
        chart_data.categories =  [k for (k, v) in tag_data]
        chart_data.add_series('Series 1', [v for (k, v) in tag_data])

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(6.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END


        self.__presentation.save(self.__file_name)

    def process_item(self, item, spider):
        # Prepare an article summary.
        text = ' '.join(item['text'])
        parser = PlaintextParser.from_string(text, Tokenizer('english'))
        summarizer = LexRankSummarizer()
        sentences = summarizer(parser.document, 4)
        summary = re.sub(r' +', ' ', ' '.join(str(s) for s in sentences))

        # Collect tags statistics.
        for tag in item['tags']:
            spider.crawler.stats.inc_value('tags/{t}'.format(t=tag))

        blank_slide_layout = self.__presentation.slide_layouts[5]
        slide = self.__presentation.slides.add_slide(blank_slide_layout)
        slide.shapes.title.text = ''.join(item['title'])
        slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(20)

        left, top, width, height = [Inches(i) for i in [1, 1.5, 8, 6]]
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        tf.text = summary
        tf.paragraphs[0].runs[0].font.size = Pt(10)
        tf.margin_left = 0
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        left, top, width, height = [Inches(i) for i in [1, 7, 8, 2]]
        tags_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tags_box.text_frame
        tags = item['tags']
        tf.text = 'URL: {u}'.format(u=item['url'])
        tf.add_paragraph().text = ', '.join(tags)
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        for paragraph in tf.paragraphs:
            paragraph.runs[0].font.size = Pt(8)

        image_path = item['images'][0]['path']
        slide.shapes.add_picture(
            os.path.join(
                spider.settings.attributes['IMAGES_STORE'].value,
                image_path),
            Inches(1), Inches(4), height=Inches(2))

        return item
